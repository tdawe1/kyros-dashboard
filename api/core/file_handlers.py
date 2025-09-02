"""
File Handlers Service

This module provides standardized file handling utilities for various file formats.
It includes support for docx, pptx, xlsx, and other common formats.
"""

import logging
from typing import Dict, Any, Optional, List, Union
from pathlib import Path

logger = logging.getLogger(__name__)


class FileHandlerError(Exception):
    """Custom exception for file handling errors."""

    pass


class FileHandler:
    """
    Base class for file handlers.
    """

    def __init__(self):
        self.supported_extensions = []
        self.max_file_size = 10 * 1024 * 1024  # 10MB default

    def can_handle(self, file_path: Union[str, Path]) -> bool:
        """
        Check if this handler can process the given file.

        Args:
            file_path: Path to the file.

        Returns:
            True if the handler can process the file.
        """
        file_path = Path(file_path)
        return file_path.suffix.lower() in self.supported_extensions

    def validate_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Validate a file before processing.

        Args:
            file_path: Path to the file.

        Returns:
            Dictionary with validation results.

        Raises:
            FileHandlerError: If validation fails.
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileHandlerError(f"File not found: {file_path}")

        if not file_path.is_file():
            raise FileHandlerError(f"Path is not a file: {file_path}")

        file_size = file_path.stat().st_size
        if file_size > self.max_file_size:
            raise FileHandlerError(
                f"File too large: {file_size} bytes (max: {self.max_file_size} bytes)"
            )

        return {
            "valid": True,
            "file_size": file_size,
            "file_path": str(file_path),
        }

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """
        Extract text content from a file.

        Args:
            file_path: Path to the file.

        Returns:
            Extracted text content.

        Raises:
            FileHandlerError: If extraction fails.
        """
        raise NotImplementedError("Subclasses must implement extract_text")

    def get_metadata(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Get metadata from a file.

        Args:
            file_path: Path to the file.

        Returns:
            Dictionary with file metadata.
        """
        file_path = Path(file_path)
        stat = file_path.stat()

        return {
            "filename": file_path.name,
            "size": stat.st_size,
            "created": stat.st_ctime,
            "modified": stat.st_mtime,
            "extension": file_path.suffix.lower(),
        }


class TextFileHandler(FileHandler):
    """
    Handler for plain text files.
    """

    def __init__(self):
        super().__init__()
        self.supported_extensions = [".txt", ".md", ".csv"]

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """Extract text from plain text files."""
        self.validate_file(file_path)

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()


class DocumentHandler(FileHandler):
    """
    Handler for Microsoft Word documents.
    """

    def __init__(self):
        super().__init__()
        self.supported_extensions = [".docx", ".doc"]

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """Extract text from Word documents."""
        self.validate_file(file_path)

        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except ImportError:
            raise FileHandlerError(
                "python-docx library not installed. Install with: pip install python-docx"
            )
        except Exception as e:
            raise FileHandlerError(f"Failed to extract text from document: {str(e)}")


class PresentationHandler(FileHandler):
    """
    Handler for Microsoft PowerPoint presentations.
    """

    def __init__(self):
        super().__init__()
        self.supported_extensions = [".pptx", ".ppt"]

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """Extract text from PowerPoint presentations."""
        self.validate_file(file_path)

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)

            return "\n".join(text_content)
        except ImportError:
            raise FileHandlerError(
                "python-pptx library not installed. Install with: pip install python-pptx"
            )
        except Exception as e:
            raise FileHandlerError(
                f"Failed to extract text from presentation: {str(e)}"
            )


class SpreadsheetHandler(FileHandler):
    """
    Handler for Microsoft Excel spreadsheets.
    """

    def __init__(self):
        super().__init__()
        self.supported_extensions = [".xlsx", ".xls"]

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """Extract text from Excel spreadsheets."""
        self.validate_file(file_path)

        try:
            import pandas as pd

            df = pd.read_excel(file_path)
            return df.to_string()
        except ImportError:
            raise FileHandlerError(
                "pandas library not installed. Install with: pip install pandas openpyxl"
            )
        except Exception as e:
            raise FileHandlerError(f"Failed to extract text from spreadsheet: {str(e)}")


class FileHandlerManager:
    """
    Manager for file handlers that can process various file types.
    """

    def __init__(self):
        self.handlers = [
            TextFileHandler(),
            DocumentHandler(),
            PresentationHandler(),
            SpreadsheetHandler(),
        ]

    def get_handler(self, file_path: Union[str, Path]) -> Optional[FileHandler]:
        """
        Get the appropriate handler for a file.

        Args:
            file_path: Path to the file.

        Returns:
            FileHandler instance or None if no handler supports the file.
        """
        for handler in self.handlers:
            if handler.can_handle(file_path):
                return handler
        return None

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """
        Extract text from a file using the appropriate handler.

        Args:
            file_path: Path to the file.

        Returns:
            Extracted text content.

        Raises:
            FileHandlerError: If no handler supports the file or extraction fails.
        """
        handler = self.get_handler(file_path)
        if not handler:
            raise FileHandlerError(f"No handler available for file: {file_path}")

        return handler.extract_text(file_path)

    def get_file_info(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Get information about a file.

        Args:
            file_path: Path to the file.

        Returns:
            Dictionary with file information.
        """
        handler = self.get_handler(file_path)
        if not handler:
            return {"error": f"No handler available for file: {file_path}"}

        try:
            validation = handler.validate_file(file_path)
            metadata = handler.get_metadata(file_path)
            return {
                "valid": True,
                "handler": handler.__class__.__name__,
                "metadata": metadata,
                "validation": validation,
            }
        except FileHandlerError as e:
            return {
                "valid": False,
                "error": str(e),
                "handler": handler.__class__.__name__,
            }

    def get_supported_extensions(self) -> List[str]:
        """
        Get list of all supported file extensions.

        Returns:
            List of supported file extensions.
        """
        extensions = []
        for handler in self.handlers:
            extensions.extend(handler.supported_extensions)
        return sorted(set(extensions))


# Global file handler manager
_file_handler_manager = None


def get_file_handler_manager() -> FileHandlerManager:
    """
    Get the global file handler manager instance.

    Returns:
        FileHandlerManager instance.
    """
    global _file_handler_manager
    if _file_handler_manager is None:
        _file_handler_manager = FileHandlerManager()
    return _file_handler_manager


def extract_text_from_file(file_path: Union[str, Path]) -> str:
    """
    Convenience function to extract text from a file.

    Args:
        file_path: Path to the file.

    Returns:
        Extracted text content.
    """
    manager = get_file_handler_manager()
    return manager.extract_text(file_path)


def get_file_info(file_path: Union[str, Path]) -> Dict[str, Any]:
    """
    Convenience function to get file information.

    Args:
        file_path: Path to the file.

    Returns:
        Dictionary with file information.
    """
    manager = get_file_handler_manager()
    return manager.get_file_info(file_path)
